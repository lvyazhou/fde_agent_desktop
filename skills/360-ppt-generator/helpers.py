# -*- coding: utf-8 -*-
"""
360 品牌 PPT Helper 函数库
=========================
复制此文件到你的脚本目录，或直接 import 使用。
所有 helper 均基于 python-pptx，配合 PPT模版.pptx 使用。

Usage:
    from helpers import *
    prs = Presentation("PPT模版.pptx")
    slide = cslide(prs)
    ptitle(slide, "我的标题")
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os
import sys
import glob


# ============ Design Tokens ============
C_BRAND     = RGBColor(0x07, 0x21, 0xA8)  # 360 品牌深蓝
C_PRIMARY   = RGBColor(0x0D, 0x47, 0xA1)  # 主色蓝
C_PRIMARY_L = RGBColor(0x42, 0xA5, 0xF5)  # 浅蓝
C_ACCENT    = RGBColor(0x00, 0xC8, 0x53)  # 强调绿
C_WARM      = RGBColor(0xFF, 0x6D, 0x00)  # 暖橙
C_PURPLE    = RGBColor(0x7C, 0x4D, 0xFF)  # 紫色
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)  # 白色
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)  # 深灰/正文色
C_GRAY      = RGBColor(0x6B, 0x7B, 0x8D)  # 次要文字
C_LGRAY     = RGBColor(0xE8, 0xEC, 0xF0)  # 浅灰背景
C_BLUE_BG   = RGBColor(0xE3, 0xF2, 0xFD)  # 蓝色浅底
C_RED       = RGBColor(0xEF, 0x53, 0x50)  # 警示红

FONT_T = "Microsoft YaHei UI"   # 标题字体
FONT_B = "Microsoft YaHei"      # 正文字体

# 模板关键索引
_MASTER2_IDX = 2       # Master 2 = 品牌装饰 master
_LY_BLANK_IDX = 6     # Master 2 下的"空白" layout


# ============ 结构函数 ============

def _find_project_root(start):
    """从 start 目录向上递归查找包含 docs/ 的项目根目录（最多向上 10 层）。"""
    cur = start
    for _ in range(10):
        if os.path.isdir(os.path.join(cur, "docs")):
            return cur
        parent = os.path.dirname(cur)
        if parent == cur:
            break
        cur = parent
    return start


def find_template(project_root=None):
    """查找 PPT模版.pptx 模板文件。

    搜索顺序：
      1. 项目 docs/ 目录下（标准位置：docs/01.../13-代理商运营/PPT模版.pptx）
      2. 项目 scripts/ 目录下
      3. 本技能目录（.roo/skills/360-ppt-generator/PPT模版.pptx）作为 fallback

    Args:
        project_root: 项目根目录。如果为 None，则从当前文件位置向上查找。

    Returns:
        模板文件的绝对路径

    Raises:
        FileNotFoundError: 找不到模板文件
    """
    if project_root is None:
        project_root = _find_project_root(
            os.path.dirname(os.path.abspath(__file__))
        )

    # 1. 搜索 docs/ 目录
    tpl_files = glob.glob(os.path.join(project_root, "docs", "**", "PPT*.pptx"), recursive=True)
    if not tpl_files:
        # 2. 搜索 scripts/ 目录
        tpl_files = glob.glob(os.path.join(project_root, "scripts", "PPT*.*"))
    if not tpl_files:
        # 3. Fallback: 技能目录自身或 brand_assets/ 子目录
        skill_dir = os.path.dirname(os.path.abspath(__file__))
        for candidate in [
            os.path.join(skill_dir, "PPT模版.pptx"),
            os.path.join(skill_dir, "brand_assets", "PPT模版.pptx"),
        ]:
            if os.path.exists(candidate):
                return candidate
        raise FileNotFoundError(
            f"Cannot find PPT template (PPT*.pptx). "
            f"Searched under: {project_root}/docs/, {project_root}/scripts/, "
            f"and {skill_dir}/ + {skill_dir}/brand_assets/"
        )
    # 优先选择 13-代理商运营 目录下的
    return next((f for f in tpl_files if "13-" in f), tpl_files[0])


def cslide(prs):
    """创建新的内容页 — 使用 Master 2 的空白 layout（带品牌装饰）。"""
    layout = prs.slide_masters[_MASTER2_IDX].slide_layouts[_LY_BLANK_IDX]
    return prs.slides.add_slide(layout)


def _move_slide_to_end(prs, idx):
    """将 prs.slides 中索引为 idx 的 slide 移动到末尾。"""
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[idx]
    sldIdLst.remove(sldId)
    sldIdLst.append(sldId)


# ============ 基础绘图 ============

def rect(slide, left, top, w, h, fill):
    """绘制圆角矩形，常用作卡片背景。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def flat_rect(slide, left, top, w, h, fill):
    """绘制直角矩形，常用作色带/背景块。"""
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def ln(slide, left, top, w, color, thick=Pt(3)):
    """绘制水平装饰线。"""
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


# ============ 文本组件 ============

def txt(slide, left, top, w, h, text, sz=18, bold=False,
        color=None, align=PP_ALIGN.LEFT, font=None):
    """添加单行/段落文本框。
    
    Args:
        slide: 目标 slide
        left, top, w, h: 位置和尺寸
        text: 文本内容（不可含 surrogate pair emoji）
        sz: 字号（磅）
        bold: 是否加粗
        color: RGBColor，默认 C_DARK
        align: 对齐方式，默认左对齐
        font: 字体名，默认 FONT_B
    """
    if color is None:
        color = C_DARK
    if font is None:
        font = FONT_B
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def multi_txt(slide, left, top, w, h, lines, sz=13, color=None, spacing=Pt(6)):
    """添加多行文本（每行一个 paragraph）。"""
    if color is None:
        color = C_DARK
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT_B
        p.space_before = spacing if i > 0 else Pt(0)
    return tb


def bullet_txt(slide, left, top, w, h, items, sz=12, color=None):
    """添加圆点列表。"""
    if color is None:
        color = C_DARK
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\u2022  " + item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT_B
        p.space_before = Pt(4) if i > 0 else Pt(0)
    return tb


# ============ 版式组件 ============

def ptitle(slide, text, subtitle=None):
    """添加页面标题 + 绿色装饰线。
    
    位置固定：左上角 (0.6", 0.25")，宽度 9"。
    """
    txt(slide, Inches(0.6), Inches(0.25), Inches(9), Inches(0.6),
        text, sz=24, bold=True, color=C_BRAND, font=FONT_T)
    ln(slide, Inches(0.6), Inches(0.82), Inches(1.5), C_ACCENT, thick=Pt(4))
    if subtitle:
        txt(slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
            subtitle, sz=12, color=C_GRAY)


def card(slide, x, y, w, h, title, body, accent=C_BRAND):
    """左边框卡片：白色背景 + 彩色左竖条 + 标题 + 正文。"""
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, Pt(5), h, accent)
    txt(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
        title, sz=13, bold=True, color=C_DARK, font=FONT_T)
    txt(slide, x + Inches(0.25), y + Inches(0.5), w - Inches(0.3), h - Inches(0.55),
        body, sz=11, color=C_GRAY)


def num_card(slide, x, y, number, label, color=C_BRAND):
    """数字指标卡片：大数字 + 说明标签。"""
    rect(slide, x, y, Inches(2.3), Inches(1.5), C_WHITE)
    txt(slide, x, y + Inches(0.15), Inches(2.3), Inches(0.7),
        number, sz=26, bold=True, color=color, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.9), Inches(2.3), Inches(0.4),
        label, sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def icon_card(slide, x, y, w, h, icon, title, desc, accent=C_BRAND):
    """图标卡片：圆形图标 + 标题 + 描述。
    
    Args:
        icon: 圆圈内显示的文字（1-2字符，如 "01"、"A"、"★"）
              注意：不可使用 surrogate pair emoji
    """
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, Pt(4), h, accent)
    # icon circle
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.15),
        Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    txt(slide, x + Inches(0.2), y + Inches(0.18), Inches(0.45), Inches(0.4),
        icon, sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.75), y + Inches(0.15), w - Inches(1.0), Inches(0.3),
        title, sz=12, bold=True, color=C_DARK, font=FONT_T)
    txt(slide, x + Inches(0.75), y + Inches(0.48), w - Inches(1.0), h - Inches(0.55),
        desc, sz=10, color=C_GRAY)


def table_slide(slide, headers, rows, x=Inches(0.6), y=Inches(1.5), w=Inches(11.0)):
    """添加品牌样式数据表格。
    
    表头：深蓝底白字
    数据行：白色/浅蓝交替
    建议每表 ≤ 8 行，超过则分页。
    默认 y=1.5" 以避免与 ptitle 重叠。
    """
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, cols, x, y, w, Inches(0.4 * (len(rows) + 1)))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = FONT_B
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_BRAND
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = C_DARK
                p.font.name = FONT_B
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if r % 2 == 0 else C_BLUE_BG
    return tbl_shape


# ============ 页面辅助 ============

def _replace_shape_text(shape, lines):
    """清空 shape 的所有段落，写入 lines = [(text, sz, bold, color, align), ...]。
    
    用于替换模板封面/结尾页已有 shape 的文字内容，避免新增 shape 导致重叠。
    
    Args:
        shape: 带有 text_frame 的 shape 对象
        lines: list of (text, font_size_pt, bold, RGBColor, PP_ALIGN)
    """
    tf = shape.text_frame
    tf.clear()
    for i, (text, sz, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_T
        p.alignment = align


def footer_note(slide, text):
    """在页面底部添加总结性结论（分隔线 + 一行蓝色文字）。
    
    位置固定在 y=6.5"，宽度 11"。每个内容页 make_xxx() 的最后一行调用。
    章节分隔页不需要调用。
    
    Args:
        slide: 目标 slide
        text: 结论文字（15-40 字为宜），会自动加 ▶ 前缀
    """
    ln(slide, Inches(0.5), Inches(6.5), Inches(11.0), C_LGRAY, Pt(1))
    txt(slide, Inches(0.5), Inches(6.55), Inches(11.0), Inches(0.5),
        f"\u25b6 {text}", sz=11, bold=False, color=C_PRIMARY, font=FONT_B)


def add_notes(slide, text):
    """为 slide 添加解说词（Speaker Notes）。

    每一页 slide 都必须调用此函数写入演讲者提词。
    解说词在演示模式下仅演讲者可见，观众看不到。

    Args:
        slide: pptx Slide 对象
        text: 解说词文本（支持 \\n 换行，50-200 字为宜）
    """
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def num_card_large(slide, x, y, number, label, color=C_BRAND):
    """加大版数字指标卡片（3.5"×2.2"），适用于汇报类 PPT。

    相比 num_card (2.3"×1.5")，字号更大、带顶部彩色装饰条。
    """
    w, h = Inches(3.5), Inches(2.2)
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, w, Pt(6), color)
    txt(slide, x, y + Inches(0.35), w, Inches(0.9),
        number, sz=36, bold=True, color=color, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(1.4), w, Inches(0.4),
        label, sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ============ 图片相关 ============

def add_image(slide, image_path, left, top, width=None, height=None):
    """在 slide 上插入图片。

    注意层级：后添加的 shape 在上层。若需图片做背景，先添加图片再添加文字。

    Args:
        slide: pptx Slide 对象
        image_path: 图片文件路径（绝对或相对）
        left, top: 位置（Inches/Emu 值）
        width, height: 尺寸。若都为 None 则保持原始大小。
                       若只传一个，另一个按比例缩放。

    Returns:
        添加的 Picture shape 对象
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image not found: {image_path}")
    return slide.shapes.add_picture(image_path, left, top, width, height)


def add_image_centered(slide, image_path, top, max_width=Inches(10), max_height=Inches(5)):
    """在 slide 上水平居中插入图片，自动缩放不超过最大宽高。

    Args:
        slide: pptx Slide 对象
        image_path: 图片文件路径
        top: 顶部 Y 坐标
        max_width: 最大宽度（默认 10"）
        max_height: 最大高度（默认 5"）

    Returns:
        添加的 Picture shape 对象
    """
    from PIL import Image
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Image not found: {image_path}")

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    # 计算缩放比例
    ratio_w = max_width / Inches(img_w / 96)  # 假设 96 DPI
    ratio_h = max_height / Inches(img_h / 96)
    ratio = min(ratio_w, ratio_h, 1.0)  # 不放大

    final_w = int(Inches(img_w / 96) * ratio)
    final_h = int(Inches(img_h / 96) * ratio)

    # 水平居中 (slide 宽度 12.2")
    left = int((Inches(12.2) - final_w) / 2)
    return slide.shapes.add_picture(image_path, left, top, final_w, final_h)


def add_bg_image(slide, image_path):
    """将图片铺满 slide 作为背景（12.2" x 6.88"），文字叠在上层保持可编辑。

    使用方法：
        slide = cslide(prs)
        add_bg_image(slide, "bg_tech_blue.png")  # 先铺背景
        txt(slide, ...)                           # 再叠文字（可编辑）

    Args:
        slide: pptx Slide 对象
        image_path: 背景图片路径（推荐 1920x1088 或 1792x1024）

    Returns:
        添加的 Picture shape 对象
    """
    if not os.path.exists(image_path):
        raise FileNotFoundError(f"Background image not found: {image_path}")
    # 全页铺满：12.2" x 6.88" (标准 16:9 slide)
    pic = slide.shapes.add_picture(image_path, Inches(0), Inches(0),
                                   Inches(12.2), Inches(6.88))
    # 将背景图移到最底层（z-order 最小）
    sp = pic._element
    sp.getparent().remove(sp)
    slide.shapes._spTree.insert(2, sp)  # index 2 = 在 spPr 之后，其余 shape 之前
    return pic


def _get_gen_script_path():
    """获取内嵌的 generate_images.py 脚本路径（本技能目录下）。"""
    skill_dir = os.path.dirname(os.path.abspath(__file__))
    local_copy = os.path.join(skill_dir, "generate_images.py")
    if os.path.exists(local_copy):
        return local_copy
    raise FileNotFoundError(
        f"Cannot find generate_images.py at: {local_copy}\n"
        f"Please ensure generate_images.py is in the skill directory."
    )


def generate_image(prompt, size="1024x1024", output_name="generated.png",
                   output_dir=None, quality="medium", timeout=300):
    """调用 image-generator skill 生成 AI 图片。

    需要系统已安装 image-generator skill。
    生成的图片保存到指定目录，供 add_image() 使用。

    Args:
        prompt: 图片描述提示词
        size: 尺寸（必须为 16 的倍数，如 1024x1024, 1792x1024）
        output_name: 输出文件名
        output_dir: 输出目录（默认 scripts/assets/）
        quality: 质量 low/medium/high
        timeout: 超时秒数（high 质量建议 300s）

    Returns:
        生成图片的完整路径

    Raises:
        RuntimeError: 图片生成失败
    """
    import subprocess

    if output_dir is None:
        output_dir = os.path.join(
            _find_project_root(os.path.dirname(os.path.abspath(__file__))),
            "scripts", "assets"
        )
    os.makedirs(output_dir, exist_ok=True)

    gen_script = _get_gen_script_path()

    cmd = [
        sys.executable, gen_script,
        "-p", prompt,
        "-s", size,
        "-q", quality,
        "-o", output_name,
        "-d", output_dir,
    ]

    # Windows 环境变量修复：避免 GBK 编码崩溃
    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"

    # 使用 bytes 模式 + 手动 decode，彻底解决 Windows 编码问题
    result = subprocess.run(
        cmd,
        stdout=subprocess.PIPE,
        stderr=subprocess.PIPE,
        timeout=timeout,
        env=env,
    )
    stdout = result.stdout.decode("utf-8", errors="replace") if result.stdout else ""
    stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""

    if result.returncode != 0:
        raise RuntimeError(
            f"Image generation failed (rc={result.returncode}): {stderr or stdout}"
        )

    output_path = os.path.join(output_dir, output_name)
    if not os.path.exists(output_path):
        raise RuntimeError(f"Expected output not found: {output_path}")
    return output_path


# PPT 配图常用尺寸常量（全部为 16 的倍数）
IMG_SIZE_FULL_BG = "1920x1088"       # 全页背景 12"x6.8"
IMG_SIZE_HALF = "960x1088"           # 半页图 5"x6.8"
IMG_SIZE_ICON = "512x512"            # 图标 ~1.5"x1.5"
IMG_SIZE_BANNER = "1792x448"         # 横幅/Banner
IMG_SIZE_SQUARE = "1024x1024"        # 正方形
IMG_SIZE_WIDE = "1792x1024"          # 宽屏 16:9


# ============ 常用 Slide Builder 模板 ============

def make_section(prs, num, title, subtitle=""):
    """章节分隔页：全色背景色带 + 大编号 + 标题。"""
    slide = cslide(prs)
    flat_rect(slide, Inches(0), Inches(2.0), Inches(12.2), Inches(2.5), C_BRAND)
    txt(slide, Inches(1.0), Inches(2.2), Inches(2), Inches(1),
        num, sz=48, bold=True, color=C_WHITE, font=FONT_T)
    txt(slide, Inches(1.0), Inches(3.0), Inches(10), Inches(0.6),
        title, sz=28, bold=True, color=C_WHITE, font=FONT_T)
    if subtitle:
        txt(slide, Inches(1.0), Inches(3.7), Inches(10), Inches(0.4),
            subtitle, sz=14, color=C_PRIMARY_L)
    return slide


# ============ 海报模式 (Poster Mode) Helpers ============

# 海报模式尺寸常量
POSTER_WIDTH = Inches(13.333)    # 16:9 标准宽
POSTER_HEIGHT = Inches(7.5)      # 16:9 标准高
POSTER_IMG_SIZE = "1792x1024"    # 图片尺寸（16:9, 16 的倍数）

# 海报模式统一 STYLE 前缀（白底 + 蓝角装饰 + 360安全云 Logo 水印）
POSTER_STYLE = (
    "Professional Chinese business presentation slide poster, "
    "CLEAN WHITE BACKGROUND with subtle blue geometric corner decorations "
    "(top-left and bottom-right corners), "
    "small '360安全云' brand watermark logo in the top-left corner "
    "(green cloud icon + '360安全云' text), "
    "modern flat design, corporate tech style with deep blue (#0721A8) text "
    "and green (#00C853) accents, "
    "clear Chinese text typography with bold titles and organized bullet points, "
    "16:9 widescreen aspect ratio. "
    "The text must be rendered clearly and accurately in Simplified Chinese. "
)


def poster_section_prompt(num, title, subtitle=""):
    """生成章节分隔页的 prompt（海报模式专用）。

    章节分隔页使用深蓝渐变背景（OVERRIDE 默认白底 STYLE），
    包含 360安全云 Logo + 大编号 + 标题 + 副标题。

    Args:
        num: 章节编号，如 "01", "02"
        title: 章节标题，如 "平台定位与价值"
        subtitle: 章节副标题，如 "为什么需要智能体协作平台"

    Returns:
        完整的 prompt 字符串（无需再拼接 STYLE）

    Example:
        >>> prompt = poster_section_prompt("01", "平台定位与价值", "为什么需要智能体协作平台")
        >>> generate_poster(prompt, "section_1")
    """
    sub_line = f"Subtitle: '{subtitle}'. " if subtitle else ""
    return (
        f"Section divider slide for Chapter {num}. "
        "OVERRIDE STYLE: deep blue gradient background with subtle tech circuit patterns. "
        "Top-left: '360安全云' logo (green cloud icon + '360安全云' text in white). "
        f"Large centered chapter number: '{num}' in bold white semi-transparent large font. "
        f"Below number: '{title}' in bold white title font. "
        f"{sub_line}"
        "Minimalist, clean, professional section divider page. "
        "This is a SECTION DIVIDER - use deep blue gradient NOT white background."
    )


def poster_cover_prompt(title, subtitle, bottom_line=""):
    """生成封面页的 prompt（海报模式专用）。

    封面页使用深蓝渐变背景（OVERRIDE 默认白底 STYLE），
    包含 360安全云 Logo 居中 + 主标题 + 副标题 + 底部信息行。

    Args:
        title: 主标题，如 "360安全云智能体安全协作平台"
        subtitle: 副标题，如 "代理商赋能培训专场"
        bottom_line: 底部信息行，如 "从入门到实战交付 · 2025年6月"

    Returns:
        完整的 prompt 字符串
    """
    bottom = f"Bottom line: '{bottom_line}'. " if bottom_line else ""
    return (
        "Title slide design for a corporate training presentation. "
        "OVERRIDE STYLE: deep blue gradient background with subtle tech circuit patterns and light particles. "
        "Top center: '360安全云' logo prominently displayed (green cloud icon + '360安全云' text in white). "
        f"Large centered title text: '{title}' in bold white font. "
        f"Subtitle below: '{subtitle}'. "
        f"{bottom}"
        "This is a COVER slide - use deep blue gradient NOT white background."
    )


def poster_closing_prompt(title="感谢聆听", subtitle="", team_info=""):
    """生成结尾页的 prompt（海报模式专用）。

    结尾页使用深蓝渐变背景，包含 360安全云 Logo + 感谢语。

    Args:
        title: 主文字，默认 "感谢聆听"
        subtitle: 副文字
        team_info: 团队/联系信息

    Returns:
        完整的 prompt 字符串
    """
    sub = f"Below: '{subtitle}'. " if subtitle else ""
    team = f"Team info: '{team_info}'. " if team_info else ""
    return (
        "Closing/thank-you slide. "
        "OVERRIDE STYLE: deep blue gradient background with subtle tech circuit patterns. "
        "Top-left: '360安全云' logo (green cloud icon + '360安全云' text in white). "
        f"Large centered text: '{title}' in bold white font. "
        f"{sub}{team}"
        "This is a CLOSING slide - use deep blue gradient NOT white background."
    )


def generate_poster(prompt, key, output_dir=None, quality="high",
                    size=POSTER_IMG_SIZE, timeout=300, regen=False):
    """生成单张海报图片（海报模式专用）。

    与 generate_image() 的区别：
    - 默认 quality="high"（确保中文渲染清晰）
    - 默认 size="1792x1024"（标准海报尺寸）
    - 默认 timeout=300（high 质量需要更长时间）
    - 支持 regen 参数控制是否跳过已有文件

    Args:
        prompt: 完整 prompt（含 STYLE 前缀）
        key: 文件名标识（不含扩展名）
        output_dir: 输出目录（默认 scripts/assets_posters/）
        quality: 质量 low/medium/high
        size: 图片尺寸
        timeout: 超时秒数
        regen: 是否强制重生成

    Returns:
        生成图片的完整路径

    Raises:
        RuntimeError: 图片生成失败
    """
    if output_dir is None:
        output_dir = os.path.join(
            _find_project_root(os.path.dirname(os.path.abspath(__file__))),
            "scripts", "assets_posters"
        )
    os.makedirs(output_dir, exist_ok=True)

    output_name = f"{key}.png"
    output_path = os.path.join(output_dir, output_name)

    if os.path.exists(output_path) and not regen:
        return output_path

    return generate_image(
        prompt=prompt,
        size=size,
        output_name=output_name,
        output_dir=output_dir,
        quality=quality,
        timeout=timeout,
    )


def assemble_poster_ppt(image_dir, slide_order, output_path):
    """将海报图片组装为 PPT（海报模式专用）。

    创建空白 16:9 演示文稿，每页全屏铺满一张海报图片。

    Args:
        image_dir: 海报图片所在目录
        slide_order: 页面顺序列表（key 列表，对应 {key}.png）
        output_path: PPT 输出文件路径

    Returns:
        (added, skipped) 元组：成功添加数和跳过数
    """
    prs = Presentation()
    prs.slide_width = POSTER_WIDTH
    prs.slide_height = POSTER_HEIGHT
    blank_layout = prs.slide_layouts[6]

    added, skipped = 0, 0
    for key in slide_order:
        img_path = os.path.join(image_dir, f"{key}.png")
        if not os.path.exists(img_path):
            skipped += 1
            continue

        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            img_path,
            left=Emu(0),
            top=Emu(0),
            width=POSTER_WIDTH,
            height=POSTER_HEIGHT,
        )
        added += 1

    os.makedirs(os.path.dirname(output_path) or ".", exist_ok=True)
    prs.save(output_path)
    return added, skipped
