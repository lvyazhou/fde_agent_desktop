# -*- coding: utf-8 -*-
"""
360 品牌 PPT 生成脚本模板
========================
基于 PPT模版.pptx 生成品牌统一的演示文稿。

使用方法:
  1. 复制此文件到 scripts/ 目录
  2. 修改 CONTENT 部分的内容
  3. 运行: python scripts/your_script.py

前置依赖:
  pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os
import glob


# ============ Design Tokens ============
C_BRAND     = RGBColor(0x07, 0x21, 0xA8)
C_PRIMARY   = RGBColor(0x0D, 0x47, 0xA1)
C_PRIMARY_L = RGBColor(0x42, 0xA5, 0xF5)
C_ACCENT    = RGBColor(0x00, 0xC8, 0x53)
C_WARM      = RGBColor(0xFF, 0x6D, 0x00)
C_PURPLE    = RGBColor(0x7C, 0x4D, 0xFF)
C_WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
C_DARK      = RGBColor(0x1A, 0x1A, 0x2E)
C_GRAY      = RGBColor(0x6B, 0x7B, 0x8D)
C_LGRAY     = RGBColor(0xE8, 0xEC, 0xF0)
C_BLUE_BG   = RGBColor(0xE3, 0xF2, 0xFD)
C_RED       = RGBColor(0xEF, 0x53, 0x50)

FONT_T = "Microsoft YaHei UI"
FONT_B = "Microsoft YaHei"
_MASTER2_IDX = 2
_LY_BLANK_IDX = 6


# ============ Helpers ============

def _move_slide_to_end(prs, idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[idx]
    sldIdLst.remove(sldId)
    sldIdLst.append(sldId)


def cslide(prs):
    """创建新内容页（Master 2 空白 layout，带品牌装饰）"""
    layout = prs.slide_masters[_MASTER2_IDX].slide_layouts[_LY_BLANK_IDX]
    return prs.slides.add_slide(layout)


def rect(slide, left, top, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def flat_rect(slide, left, top, w, h, fill):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def ln(slide, left, top, w, color, thick=Pt(3)):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, thick)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def txt(slide, left, top, w, h, text, sz=18, bold=False,
        color=None, align=PP_ALIGN.LEFT, font=None):
    if color is None:
        color = C_DARK
    if font is None:
        font = FONT_B
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font
    p.alignment = align
    return tb


def multi_txt(slide, left, top, w, h, lines, sz=13, color=None, spacing=Pt(6)):
    if color is None:
        color = C_DARK
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(4)
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT_B
        p.space_before = spacing if i > 0 else Pt(0)
    return tb


def bullet_txt(slide, left, top, w, h, items, sz=12, color=None):
    if color is None:
        color = C_DARK
    tb = slide.shapes.add_textbox(left, top, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_top = Pt(2)
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "\u2022  " + item
        p.font.size = Pt(sz)
        p.font.color.rgb = color
        p.font.name = FONT_B
        p.space_before = Pt(4) if i > 0 else Pt(0)
    return tb


def ptitle(slide, text, subtitle=None):
    txt(slide, Inches(0.6), Inches(0.25), Inches(9), Inches(0.6),
        text, sz=24, bold=True, color=C_BRAND, font=FONT_T)
    ln(slide, Inches(0.6), Inches(0.82), Inches(1.5), C_ACCENT, thick=Pt(4))
    if subtitle:
        txt(slide, Inches(0.6), Inches(0.88), Inches(10), Inches(0.3),
            subtitle, sz=12, color=C_GRAY)


def card(slide, x, y, w, h, title, body, accent=C_BRAND):
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, Pt(5), h, accent)
    txt(slide, x + Inches(0.25), y + Inches(0.12), w - Inches(0.3), Inches(0.35),
        title, sz=13, bold=True, color=C_DARK, font=FONT_T)
    txt(slide, x + Inches(0.25), y + Inches(0.5), w - Inches(0.3), h - Inches(0.55),
        body, sz=11, color=C_GRAY)


def num_card(slide, x, y, number, label, color=C_BRAND):
    rect(slide, x, y, Inches(2.3), Inches(1.5), C_WHITE)
    txt(slide, x, y + Inches(0.15), Inches(2.3), Inches(0.7),
        number, sz=26, bold=True, color=color, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(0.9), Inches(2.3), Inches(0.4),
        label, sz=11, color=C_GRAY, align=PP_ALIGN.CENTER)


def icon_card(slide, x, y, w, h, icon, title, desc, accent=C_BRAND):
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, Pt(4), h, accent)
    circ = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, x + Inches(0.2), y + Inches(0.15),
        Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = accent
    circ.line.fill.background()
    txt(slide, x + Inches(0.2), y + Inches(0.18), Inches(0.45), Inches(0.4),
        icon, sz=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    txt(slide, x + Inches(0.75), y + Inches(0.15), w - Inches(1.0), Inches(0.3),
        title, sz=12, bold=True, color=C_DARK, font=FONT_T)
    txt(slide, x + Inches(0.75), y + Inches(0.48), w - Inches(1.0), h - Inches(0.55),
        desc, sz=10, color=C_GRAY)


def table_slide(slide, headers, rows, x=Inches(0.6), y=Inches(1.5), w=Inches(11.0)):
    cols = len(headers)
    tbl_shape = slide.shapes.add_table(
        len(rows) + 1, cols, x, y, w, Inches(0.4 * (len(rows) + 1)))
    tbl = tbl_shape.table
    for i, h in enumerate(headers):
        cell = tbl.cell(0, i)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = C_WHITE
            p.font.name = FONT_B
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_BRAND
    for r, row_data in enumerate(rows):
        for c, val in enumerate(row_data):
            cell = tbl.cell(r + 1, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.color.rgb = C_DARK
                p.font.name = FONT_B
            cell.fill.solid()
            cell.fill.fore_color.rgb = C_WHITE if r % 2 == 0 else C_BLUE_BG
    return tbl_shape


def make_section(prs, num, title, subtitle=""):
    """章节分隔页（不需要 footer_note）"""
    slide = cslide(prs)
    flat_rect(slide, Inches(0), Inches(2.0), Inches(12.2), Inches(2.5), C_BRAND)
    txt(slide, Inches(1.0), Inches(2.2), Inches(2), Inches(1),
        num, sz=48, bold=True, color=C_WHITE, font=FONT_T)
    txt(slide, Inches(1.0), Inches(3.0), Inches(10), Inches(0.6),
        title, sz=28, bold=True, color=C_WHITE, font=FONT_T)
    if subtitle:
        txt(slide, Inches(1.0), Inches(3.7), Inches(10), Inches(0.4),
            subtitle, sz=14, color=C_PRIMARY_L)
    return slide


# ============ 封面/结尾页 & 底部结论 ============

def _replace_shape_text(shape, lines):
    """清空 shape 文本并写入新内容（用于封面/结尾页已有 shape）。
    lines = [(text, sz, bold, color, align), ...]
    """
    tf = shape.text_frame
    tf.clear()
    for i, (text, sz, bold, color, align) in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = text
        p.font.size = Pt(sz)
        p.font.bold = bold
        p.font.color.rgb = color
        p.font.name = FONT_T
        p.alignment = align


def update_cover(prs, title, subtitle):
    """更新模板封面（idx 0）：替换已有 TextBox 文字，避免重叠。
    封面背景白色 → 必须用 C_BRAND/C_GRAY 等深色文字。
    """
    slide = prs.slides[0]
    for shape in slide.shapes:
        if shape.has_text_frame and "360" in shape.text_frame.text:
            _replace_shape_text(shape, [
                (title, 40, True, C_BRAND, PP_ALIGN.CENTER),
                ("", 14, False, C_BRAND, PP_ALIGN.CENTER),
                (subtitle, 18, False, C_GRAY, PP_ALIGN.CENTER),
            ])
            return
    # fallback
    txt(slide, Inches(1.5), Inches(2.2), Inches(9), Inches(1.2),
        title, sz=40, bold=True, color=C_BRAND, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, Inches(1.5), Inches(3.6), Inches(9), Inches(0.5),
        subtitle, sz=18, color=C_GRAY, align=PP_ALIGN.CENTER)


def update_closing(prs, title="感谢聆听", subtitle=""):
    """更新模板结尾页（idx 2，稍后被 _move_slide_to_end 移到末尾）。
    结尾页背景白色 → 必须用 C_BRAND/C_GRAY 等深色文字。
    """
    slide = prs.slides[2]
    for shape in slide.shapes:
        if shape.has_text_frame and "贡献" in shape.text_frame.text:
            _replace_shape_text(shape, [
                (title, 40, True, C_BRAND, PP_ALIGN.CENTER),
                ("", 14, False, C_BRAND, PP_ALIGN.CENTER),
                (subtitle, 16, False, C_GRAY, PP_ALIGN.CENTER),
            ])
            return
    # fallback
    txt(slide, Inches(2.0), Inches(2.5), Inches(8), Inches(0.9),
        title, sz=40, bold=True, color=C_BRAND, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, Inches(2.0), Inches(3.6), Inches(8), Inches(0.5),
        subtitle, sz=16, color=C_GRAY, align=PP_ALIGN.CENTER)


def footer_note(slide, text):
    """每个内容页 make_xxx() 最后一行调用，在 y=6.5" 添加总结结论。
    章节分隔页不需要。结论文字 15-40 字为宜。
    """
    ln(slide, Inches(0.5), Inches(6.5), Inches(11.0), C_LGRAY, Pt(1))
    txt(slide, Inches(0.5), Inches(6.55), Inches(11.0), Inches(0.5),
        f"\u25b6 {text}", sz=11, bold=False, color=C_PRIMARY, font=FONT_B)


def num_card_large(slide, x, y, number, label, color=C_BRAND):
    """加大版数字卡 (3.5"×2.2")，比默认 num_card 更醒目，适合汇报类 PPT。"""
    w, h = Inches(3.5), Inches(2.2)
    rect(slide, x, y, w, h, C_WHITE)
    flat_rect(slide, x, y, w, Pt(6), color)
    txt(slide, x, y + Inches(0.35), w, Inches(0.9),
        number, sz=36, bold=True, color=color, font=FONT_T, align=PP_ALIGN.CENTER)
    txt(slide, x, y + Inches(1.4), w, Inches(0.4),
        label, sz=12, color=C_GRAY, align=PP_ALIGN.CENTER)


# ============ 图片插入 ============

def add_image(slide, image_path, left, top, width=None, height=None):
    """在 slide 上插入图片（后添加的 shape 在上层）。

    Args:
        slide: pptx Slide 对象
        image_path: 图片路径
        left, top: 位置（Inches/Emu）
        width, height: 尺寸，None 则保持原始比例
    Returns:
        Picture shape 对象
    """
    if not os.path.exists(image_path):
        print(f"[WARN] Image not found, skipped: {image_path}")
        return None
    return slide.shapes.add_picture(image_path, left, top, width, height)


# PPT 配图常用尺寸（全部为 16 的倍数，供 image-generator 使用）
IMG_SIZE_FULL_BG = "1920x1088"   # 全页背景
IMG_SIZE_HALF = "960x1088"       # 半页图
IMG_SIZE_ICON = "512x512"        # 图标
IMG_SIZE_WIDE = "1792x1024"      # 宽屏 16:9


# ====================== SLIDE BUILDERS ======================
# 在这里编写你的 make_xxx(prs) 函数
# 每个函数负责一张 slide 的内容


def make_sample_title_cards(prs):
    """示例：3 列卡片布局"""
    slide = cslide(prs)
    ptitle(slide, "示例：卡片布局")
    cards_data = [
        ("功能一", "这是功能一的\n详细描述内容"),
        ("功能二", "这是功能二的\n详细描述内容"),
        ("功能三", "这是功能三的\n详细描述内容"),
    ]
    for i, (title, body) in enumerate(cards_data):
        card(slide, Inches(0.5) + i * Inches(3.7), Inches(1.5),
             Inches(3.5), Inches(2.0), title, body,
             accent=[C_BRAND, C_ACCENT, C_WARM][i])


def make_sample_metrics(prs):
    """示例：数字指标卡片（使用加大版 + 底部结论）"""
    slide = cslide(prs)
    ptitle(slide, "示例：数据概览")
    metrics = [
        ("99.9%", "系统可用性"),
        ("50+", "已建智能体"),
        ("200+", "技能组件"),
        ("3", "级租户体系"),
        ("6", "级角色权限"),
        ("100%", "审计覆盖"),
    ]
    for i, (num, label) in enumerate(metrics):
        row, col = i // 3, i % 3
        x = Inches(0.4) + col * Inches(3.9)
        y = Inches(1.5) + row * Inches(2.6)
        num_card_large(slide, x, y, num, label,
                       color=C_BRAND if row == 0 else C_PRIMARY)
    footer_note(slide, "结论：核心指标全面达成，系统稳定性与覆盖率持续提升。")


def make_sample_table(prs):
    """示例：数据表格 + 底部结论"""
    slide = cslide(prs)
    ptitle(slide, "示例：表格布局")
    headers = ["阶段", "任务", "负责人", "时间"]
    rows = [
        ["第1周", "环境部署", "技术团队", "2天"],
        ["第2周", "数据对接", "实施团队", "3天"],
        ["第3周", "智能体配置", "AI 工程师", "5天"],
        ["第4周", "UAT 测试", "客户方", "3天"],
        ["第5周", "正式上线", "项目经理", "2天"],
    ]
    table_slide(slide, headers, rows)
    footer_note(slide, "结论：5 周完成全部交付，关键路径在第 3 周智能体配置。")


# ====================== MAIN ======================

def _find_project_root(start):
    """从 start 目录向上查找包含 docs/ 的项目根目录（最多向上 10 层）。"""
    cur = start
    for _ in range(10):
        if os.path.isdir(os.path.join(cur, "docs")):
            return cur
        parent = os.path.dirname(cur)
        if parent == cur:
            break
        cur = parent
    return start


def main():
    base = os.path.dirname(os.path.abspath(__file__))
    # 兼容三种运行位置：scripts/、.roo/skills/360-ppt-generator/、项目根
    project_root = _find_project_root(base)

    # 1. 查找模板（docs/ → scripts/ → 技能目录自身）
    tpl_files = glob.glob(
        os.path.join(project_root, "docs", "**", "PPT*.pptx"), recursive=True)
    if not tpl_files:
        tpl_files = glob.glob(os.path.join(project_root, "scripts", "PPT*.*"))
    if not tpl_files:
        # Fallback: 技能目录自带的模板副本
        skill_dir = os.path.dirname(os.path.abspath(__file__))
        local_tpl = os.path.join(skill_dir, "PPT模版.pptx")
        if os.path.exists(local_tpl):
            tpl_files = [local_tpl]
        else:
            raise FileNotFoundError(
                f"Cannot find PPT template (PPT*.pptx) under "
                f"{project_root}/docs/, {project_root}/scripts/, or {skill_dir}/")
    tpl = next((f for f in tpl_files if "13-" in f), tpl_files[0])
    print(f"Using template: {tpl}")
    prs = Presentation(tpl)

    # 2. 更新封面 + 结尾（模板 slide idx 0 / idx 2）
    update_cover(prs, "你的标题", "副标题  │  部门  │  日期")
    update_closing(prs, "感谢聆听", "部门名称  │  日期")

    # 3. 添加内容页（每个内容页内部应调用 footer_note）
    make_section(prs, "01", "第一章", "副标题描述")
    make_sample_title_cards(prs)
    make_sample_metrics(prs)
    make_sample_table(prs)

    # 4. 将结尾页移到末尾
    _move_slide_to_end(prs, 2)

    # 5. 保存
    out = os.path.join(base, "_output.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Total slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
