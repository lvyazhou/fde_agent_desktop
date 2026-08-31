# -*- coding: utf-8 -*-
"""
混合模式 PPT 脚本模板
=====================
同一个 PPT 中组合海报页（AI 生成全页图片）和文字页（python-pptx 精确排版）。

典型用法：
- 封面 / 章节分隔 / 结尾 → 海报模式（视觉冲击力）
- 数据 / 表格 / 要点列表 → 文字模式（精确可编辑）

使用方法:
  python template_hybrid_script.py
"""
import os
import sys
import time

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# ============ 路径设置 ============
SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, SCRIPT_DIR)

ASSETS = os.path.join(SCRIPT_DIR, "assets_hybrid")
os.makedirs(ASSETS, exist_ok=True)

# ============ 品牌色 ============
C_BRAND = RGBColor(0x07, 0x21, 0xA8)   # 360 深蓝
C_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)
C_DARK = RGBColor(0x1A, 0x1A, 0x2E)
C_ACCENT = RGBColor(0x00, 0xC8, 0x53)  # 360 绿

# ============ 模板 Layout 索引 ============
_LY_BLANK_IDX = 6  # Master 2 的空白 layout（所有内容页使用）


# ============ 基础工具函数 ============

def _find_project_root(start):
    """向上查找包含 docs/ 目录的项目根。"""
    d = start
    for _ in range(10):
        if os.path.isdir(os.path.join(d, "docs")):
            return d
        parent = os.path.dirname(d)
        if parent == d:
            break
        d = parent
    return start


def find_template(project_root=None):
    """3-level 查找 PPT 模板文件。"""
    if project_root is None:
        project_root = _find_project_root(SCRIPT_DIR)

    candidates = [
        os.path.join(project_root, "docs", "01 医院智能体安全协作平台",
                     "13-代理商运营", "PPT模版.pptx"),
        os.path.join(project_root, "scripts", "PPT模版.pptx"),
        os.path.join(SCRIPT_DIR, "PPT模版.pptx"),
    ]
    for p in candidates:
        if os.path.isfile(p):
            return p
    raise FileNotFoundError(
        f"PPT模版.pptx not found in:\n" + "\n".join(f"  - {c}" for c in candidates)
    )


def cslide(prs):
    """创建基于 Master 2 空白 layout 的新 slide（带品牌装饰）。"""
    layout = prs.slide_layouts[_LY_BLANK_IDX]
    return prs.slides.add_slide(layout)


def _move_slide_to_end(prs, idx):
    """将指定索引的 slide 移到末尾。"""
    slides = prs.slides._sldIdLst
    el = slides[idx]
    slides.remove(el)
    slides.append(el)


# ============ 文字模式工具函数 ============

def ptitle(slide, text, subtitle=None):
    """添加页面标题。"""
    from pptx.util import Inches, Pt
    tf = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), Inches(10), Inches(0.6)).text_frame
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = C_BRAND
    if subtitle:
        sp = tf.add_paragraph()
        sp.text = subtitle
        sp.font.size = Pt(13)
        sp.font.color.rgb = RGBColor(0x66, 0x66, 0x66)


def txt(slide, left, top, w, h, text, sz=18, bold=False, color=None, align=None):
    """添加文字框。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(sz)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color
    if align:
        p.alignment = align
    return tf


def bullet_txt(slide, left, top, w, h, items, sz=12, color=None):
    """添加项目符号列表。"""
    from pptx.util import Pt
    tf = slide.shapes.add_textbox(left, top, w, h).text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"•  {item}"
        p.font.size = Pt(sz)
        if color:
            p.font.color.rgb = color
    return tf


def num_card(slide, x, y, number, label, color=C_BRAND):
    """数字指标卡。"""
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN
    w, h = Inches(2.5), Inches(1.8)
    shape = slide.shapes.add_shape(1, x, y, w, h)  # MSO_SHAPE.ROUNDED_RECTANGLE
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_LIGHT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = str(number)
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER
    p2 = tf.add_paragraph()
    p2.text = label
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x66, 0x66, 0x66)
    p2.alignment = PP_ALIGN.CENTER


def card(slide, x, y, w, h, title, body, accent=C_BRAND):
    """内容卡片。"""
    from pptx.util import Inches, Pt
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = C_LIGHT
    shape.line.fill.background()
    tf = shape.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = accent
    p2 = tf.add_paragraph()
    p2.text = body
    p2.font.size = Pt(11)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# ============ 海报模式工具函数 ============

# STYLE 前缀 - 白底（内容页默认）
POSTER_STYLE = (
    "Professional Chinese business presentation slide poster, "
    "CLEAN WHITE BACKGROUND with subtle blue geometric corner decorations "
    "(top-left and bottom-right corners), "
    "small '360安全云' brand watermark logo in the top-left corner "
    "(green cloud icon + '360安全云' text), "
    "modern flat design, corporate tech style with deep blue (#0721A8) text "
    "and green (#00C853) accents, "
    "clear Chinese text typography with bold titles and organized bullet points, "
    "16:9 widescreen aspect ratio. "
)

# STYLE_DARK 前缀 - 深蓝底（封面/章节/结尾）
POSTER_STYLE_DARK = (
    "Professional Chinese business presentation slide poster, "
    "DEEP BLUE GRADIENT BACKGROUND (#0721A8 to #0a0e2a), "
    "large '360安全云' brand logo prominently placed (green cloud icon + white text), "
    "modern flat design with white text and green (#00C853) accents, "
    "clear Chinese text typography with bold titles, "
    "16:9 widescreen aspect ratio. "
)

IMAGE_SIZE = "1792x1024"


def _find_gen_script():
    """查找内嵌的 generate_images.py（本技能目录下）。"""
    local = os.path.join(SCRIPT_DIR, "generate_images.py")
    if os.path.isfile(local):
        return local
    raise FileNotFoundError(
        f"generate_images.py not found at: {local}\n"
        f"Please ensure generate_images.py is in the skill directory."
    )


GEN_SCRIPT = _find_gen_script()


def generate_poster(prompt, key, output_dir=None, quality="high", size=IMAGE_SIZE):
    """生成单张海报图片。"""
    import subprocess
    if output_dir is None:
        output_dir = ASSETS
    out_path = os.path.join(output_dir, f"{key}.png")
    if os.path.exists(out_path):
        print(f"  [skip] {key}.png already exists")
        return out_path

    cmd = [
        sys.executable, GEN_SCRIPT,
        "-p", prompt,
        "-s", size,
        "-q", quality,
        "-o", f"{key}.png",
        "-d", output_dir,
    ]
    print(f"  [gen] {key} ...")
    try:
        subprocess.run(cmd, check=True, capture_output=True, text=True, timeout=180)
        print(f"  [done] {key}.png")
    except subprocess.CalledProcessError as e:
        print(f"  [ERROR] {key}: {e.stderr[:200]}")
    except subprocess.TimeoutExpired:
        print(f"  [TIMEOUT] {key}")

    return out_path


def add_image_centered(slide, image_path, top, max_width=Inches(10), max_height=Inches(5)):
    """在 slide 上水平居中插入图片，自动缩放。"""
    try:
        from PIL import Image
        img = Image.open(image_path)
        img_w, img_h = img.size
        img.close()
    except ImportError:
        # 没有 Pillow 就用固定尺寸
        slide.shapes.add_picture(image_path, Inches(0.1), top, max_width)
        return

    from pptx.util import Emu
    aspect = img_w / img_h
    # 先按宽度缩放
    width = max_width
    height = Emu(int(width / aspect))
    # 如果高度超限，按高度缩放
    if height > max_height:
        height = max_height
        width = Emu(int(height * aspect))
    # 居中
    slide_width = Inches(12.2)
    left = (slide_width - width) // 2
    slide.shapes.add_picture(image_path, left, top, width, height)


# ============ 海报页定义 ============
# 格式: (key, style + prompt)

POSTER_PAGES = [
    # ---- 封面 ----
    ("cover", POSTER_STYLE_DARK +
     "封面页。大标题: '产品介绍', 副标题: '360安全云智能体协作平台', "
     "底部一行小字: '360 数字安全集团'"),

    # ---- 章节分隔页 ----
    ("section_01", POSTER_STYLE_DARK +
     "章节分隔页。左侧超大半透明数字 '01', "
     "右侧标题: '产品概述', 副标题: '定位与价值'"),

    ("section_02", POSTER_STYLE_DARK +
     "章节分隔页。左侧超大半透明数字 '02', "
     "右侧标题: '核心能力', 副标题: '平台特性与指标'"),

    # ---- 结尾 ----
    ("closing", POSTER_STYLE_DARK +
     "结尾页。中央大字: '感谢聆听', 下方: '360安全云团队'"),
]


def gen_poster_images(only=None, regen=False):
    """批量生成海报页图片。"""
    for key, prompt in POSTER_PAGES:
        if only and key not in only:
            continue
        out = os.path.join(ASSETS, f"{key}.png")
        if os.path.exists(out) and not regen:
            print(f"  [skip] {key}.png")
            continue
        generate_poster(prompt, key)
        time.sleep(2)  # API 限流保护


# ============ 文字模式内容页 ============

def make_agenda(prs):
    """目录页（文字模式）。"""
    slide = cslide(prs)
    ptitle(slide, "目录")
    items = ["01  产品概述", "02  核心能力", "03  实施案例", "04  下一步计划"]
    bullet_txt(slide, Inches(1.0), Inches(1.8), Inches(10), Inches(4), items, sz=18)


def make_product_overview(prs):
    """产品定位页（文字模式）。"""
    slide = cslide(prs)
    ptitle(slide, "产品定位")
    card(slide, Inches(0.6), Inches(1.5), Inches(5), Inches(2),
         "SaaS 平台", "多租户 AI 智能体协作平台")
    card(slide, Inches(6.2), Inches(1.5), Inches(5), Inches(2),
         "安全合规", "等保三级 + 数据隔离")


def make_kpi_cards(prs):
    """核心指标页（文字模式）。"""
    slide = cslide(prs)
    ptitle(slide, "核心指标")
    num_card(slide, Inches(0.6), Inches(2.0), "150+", "接入医院")
    num_card(slide, Inches(3.4), Inches(2.0), "500+", "活跃智能体")
    num_card(slide, Inches(6.2), Inches(2.0), "99.9%", "系统可用性")
    num_card(slide, Inches(9.0), Inches(2.0), "50M+", "月度 Token")


# ============ 混合组装 ============

def _insert_poster_slide(prs, key):
    """插入一张海报页（全页覆盖图片）。"""
    img_path = os.path.join(ASSETS, f"{key}.png")
    if not os.path.exists(img_path):
        print(f"  [WARN] {key}.png not found, creating empty slide")
        cslide(prs)
        return
    slide = cslide(prs)
    add_image_centered(slide, img_path, Inches(0),
                       max_width=Inches(12.2), max_height=Inches(6.88))


def main():
    """混合模式主函数。"""
    import argparse
    parser = argparse.ArgumentParser(description="混合模式 PPT 生成")
    parser.add_argument("--regen", action="store_true", help="重新生成所有海报图片")
    parser.add_argument("--skip-posters", action="store_true", help="跳过海报生成（仅组装）")
    args = parser.parse_args()

    # ---- Step 1: 生成海报图片 ----
    if not args.skip_posters:
        print("=" * 50)
        print("Step 1: 生成海报页图片")
        print("=" * 50)
        gen_poster_images(regen=args.regen)

    # ---- Step 2: 加载模板 ----
    print("\n" + "=" * 50)
    print("Step 2: 加载 360 品牌模板 & 组装 PPT")
    print("=" * 50)
    tpl = find_template()
    prs = Presentation(tpl)
    print(f"  Template: {tpl}")

    # ---- Step 3: 覆盖封面 ----
    cover_img = os.path.join(ASSETS, "cover.png")
    if os.path.exists(cover_img):
        slide = prs.slides[0]
        add_image_centered(slide, cover_img, Inches(0),
                           max_width=Inches(12.2), max_height=Inches(6.88))
        print("  [slide 0] Cover → poster")

    # ---- Step 4: 按顺序添加内容页 ----
    make_agenda(prs)                        # 文字页: 目录
    print("  [slide +] Agenda → text")

    _insert_poster_slide(prs, "section_01")  # 海报页: 章节1
    print("  [slide +] Section 01 → poster")

    make_product_overview(prs)               # 文字页: 产品定位
    print("  [slide +] Product Overview → text")

    _insert_poster_slide(prs, "section_02")  # 海报页: 章节2
    print("  [slide +] Section 02 → poster")

    make_kpi_cards(prs)                      # 文字页: KPI
    print("  [slide +] KPI Cards → text")

    # ---- Step 5: 结尾页 ----
    closing_img = os.path.join(ASSETS, "closing.png")
    closing_slide = prs.slides[2]  # 模板结尾页（原始 idx 2）
    if os.path.exists(closing_img):
        add_image_centered(closing_slide, closing_img, Inches(0),
                           max_width=Inches(12.2), max_height=Inches(6.88))
    _move_slide_to_end(prs, 2)
    print("  [slide →] Closing → moved to end")

    # ---- Step 6: 保存 ----
    out = os.path.join(os.path.dirname(__file__), "output_hybrid.pptx")
    prs.save(out)
    print(f"\n{'=' * 50}")
    print(f"Done! {len(prs.slides)} slides → {out}")
    print(f"{'=' * 50}")


if __name__ == "__main__":
    main()
