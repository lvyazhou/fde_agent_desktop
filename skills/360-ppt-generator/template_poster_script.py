# -*- coding: utf-8 -*-
"""
360 PPT 海报模式 — 参考模板脚本（v4 - 浅蓝白底风格 / 真人验货确认）
=======================================================

核心风格（v4，2026-07 确认）：
  - 内容页统一「浅蓝白底渐变 + 圆角卡片」腾讯健康产品页风格（见下方 STYLE）
  - 深蓝色 (#0721A8) 标题、绿(#00C853)/蓝点缀、大留白
  - 标题放在 UPPER-CENTER，绝不占用左上角
  - 左上角必须整块留空（top 1.2" × left 3.5"）给真 logo PNG 叠加

核心机制（v3 起）：
  - STYLE 中不要求 AI 画 logo（AI 画的 logo 每页都不一样，且会串到 360 字样）
  - 组装 PPT 时用 python-pptx 把真实 logo PNG 叠加到每页固定位置
  - logo 来源：brand_assets/logo-title.png（深色，配浅底）
                brand_assets/logo-title-white.png（白色，配深底封面/结尾）

使用方法：
  - 复制此模板到 scripts/ 目录
  - 修改 STYLE、POSTERS 为实际内容
  - 运行: python scripts/your_script.py gen     # 生成图片
  - 运行: python scripts/your_script.py build   # 组装 PPT
  - 运行: python scripts/your_script.py all     # 全流程

页面类型（通过 POSTERS 中的 page_type 字段区分）：
  - "cover"   : 封面页 — 可深蓝渐变或浅蓝白底装饰背景，大 logo 居中偏上
  - "section" : 章节分隔页 — 深蓝渐变背景，中 logo 左上角
  - "content" : 内容页 — 浅蓝白底渐变（默认 STYLE），小 logo 左上角
  - "closing" : 结尾页 — 深蓝渐变或浅蓝白底，大 logo 居中

依赖：
  - python-pptx >= 0.6.21
  - Pillow (PIL)
  - image-generator skill (generate_images.py)
"""

import os
import sys
import subprocess
import time
import argparse
from pptx import Presentation
from pptx.util import Inches, Emu
from PIL import Image

# ╔══════════════════════════════════════════════════════════════╗
# ║  配置区 - 根据实际项目修改                                      ║
# ╚══════════════════════════════════════════════════════════════╝

# 输出目录
BASE_DIR = os.path.dirname(os.path.abspath(__file__))
ASSETS = os.path.join(BASE_DIR, "assets_posters")
os.makedirs(ASSETS, exist_ok=True)

# Logo 路径（真实 360安全云 logo PNG，透明背景）
SKILL_DIR = os.path.dirname(os.path.abspath(__file__))
LOGO_PATH = os.path.join(SKILL_DIR, "brand_assets", "logo-title.png")

# image-generator 脚本路径
def _find_gen_script():
    """查找内嵌的 generate_images.py（本技能目录下）。"""
    skill_dir = os.path.dirname(os.path.abspath(__file__))
    local_copy = os.path.join(skill_dir, "generate_images.py")
    if os.path.exists(local_copy):
        return local_copy
    raise FileNotFoundError(
        f"generate_images.py not found at: {local_copy}\n"
        f"Please ensure generate_images.py is in the skill directory."
    )

GEN_SCRIPT = _find_gen_script()

# PPT 尺寸
SLIDE_WIDTH = Inches(13.333)   # 16:9 标准宽
SLIDE_HEIGHT = Inches(7.5)     # 16:9 标准高

# 图片规格
IMAGE_SIZE = "1792x1024"       # 16:9, 16 的倍数
IMAGE_QUALITY = "high"         # low / medium / high
TIMEOUT_SEC = 300              # 单张超时
SLEEP_SEC = 2                  # 请求间隔

# ╔══════════════════════════════════════════════════════════════╗
# ║  Logo 叠加配置                                                 ║
# ║  三类页面的 logo 位置和尺寸                                      ║
# ╚══════════════════════════════════════════════════════════════╝

# logo-title.png 原始尺寸 3800x1160，内容区 3000x505
# 宽高比约 5.94:1

LOGO_PRESETS = {
    # 封面页：大 logo，居中偏上
    "cover": {
        "width": Inches(3.5),       # 大尺寸
        "left": Inches(4.9),        # 水平居中 (13.333 - 3.5) / 2
        "top": Inches(1.2),         # 偏上
    },
    # 章节分隔页：中 logo，左上角
    "section": {
        "width": Inches(2.2),       # 中等尺寸
        "left": Inches(0.5),        # 左上
        "top": Inches(0.3),
    },
    # 内容页：logo 左上角
    "content": {
        "width": Inches(2.2),       # 与章节页统一
        "left": Inches(0.2),
        "top": Inches(0.1),
    },
    # 结尾页：大 logo，居中
    "closing": {
        "width": Inches(3.5),
        "left": Inches(4.9),
        "top": Inches(1.5),
    },
}

# ╔══════════════════════════════════════════════════════════════╗
# ║  统一视觉风格（STYLE 前缀）                                     ║
# ║  *** 不再包含 logo 描述 — logo 由 python-pptx 叠加 ***         ║
# ╚══════════════════════════════════════════════════════════════╝

# 内容页默认风格：浅蓝白底渐变 + 卡片式布局（腾讯健康产品页风格，无 logo）
# —— 2026-07 经真人验货确认的风格，务必保留 LOGO SPACE 与 upper-center 标题规则
STYLE = (
    "Professional Chinese business presentation slide, premium Tencent-healthcare "
    "product-page aesthetic. CLEAN LIGHT-BLUE-TO-WHITE gradient background (very light, "
    "airy, lots of whitespace), with subtle thin blue geometric decorations only in the "
    "TOP-RIGHT and BOTTOM-RIGHT corners. "
    "CRITICAL LAYOUT RULE - LOGO SPACE: The ENTIRE top-left corner area (the top 1.2 inches "
    "height, left 3.5 inches width) MUST be COMPLETELY EMPTY - pure background color only, "
    "absolutely NO text, NO graphics, NO decorations, NO title in that zone. A real company "
    "logo PNG will be placed there afterward. Put the page title in the UPPER-CENTER, "
    "starting well clear of that top-left logo zone. "
    "STRICTLY NO LOGO, NO BRAND MARK, NO '360' TEXT anywhere - a real logo PNG will be "
    "overlaid later. "
    "Modern flat design, rounded cards with soft shadows, deep blue (#0721A8) titles, "
    "green (#00C853) and blue accents, clear accurate Simplified Chinese typography with "
    "bold titles and well-organized content, generous spacing. 16:9 widescreen. "
    "FONT: render ALL Chinese text in a clean modern sans-serif GOTHIC typeface like "
    "Microsoft YaHei / Source Han Sans - no serif, no decorative or handwriting fonts, "
    "consistent across the whole slide. "
    "FONT-SIZE HIERARCHY (keep this clear visual scale): page title largest and bold "
    "(about 40pt), section/card headings medium bold (about 22-24pt), body text regular "
    "(about 18pt), secondary notes / captions smaller (about 14pt). "
    "The text must be rendered clearly and accurately in Simplified Chinese. "
)

# 深蓝背景风格：用于封面 / 章节分隔 / 结尾页
# 注意：封面/结尾页通常直接在 prompt 里写"浅蓝白底装饰背景"，只有需要深蓝时才叠加 STYLE_DARK
STYLE_DARK = (
    "OVERRIDE STYLE: Use a DEEP BLUE gradient background (#0B1D4A to #0721A8) "
    "instead of light background. "
    "DO NOT draw any logo or brand mark — a real logo PNG will be overlaid later. "
    "All text in WHITE or light blue (#E5E9FA). "
    "Decorative: subtle tech grid pattern, geometric accents, glowing particle effects. "
    "Keep the top-center/top-left area clean for a logo overlay. "
)

# ╔══════════════════════════════════════════════════════════════╗
# ║  海报内容定义                                                   ║
# ║  Tuple 格式: (key, page_type, prompt)                          ║
# ║  page_type: "cover" | "section" | "content" | "closing"       ║
# ╚══════════════════════════════════════════════════════════════╝

POSTERS = [
    # ═══════ 封面（纯装饰背景，文字由 python-pptx 叠加） ═══════
    ("cover", "cover", STYLE + STYLE_DARK +
        "Design a PURELY DECORATIVE background for a cover/title slide.\n"
        "STRICTLY NO TEXT, NO CHARACTERS, NO NUMBERS, NO LETTERS, NO WORDS of any language.\n"
        "Leave top-center area clean for logo overlay.\n"
        "Leave the center area (middle 60% of the image) relatively clean "
        "for text to be overlaid later by code — avoid placing major visual elements there.\n"
        "Decorative elements on edges only: subtle radial glow, geometric frame, "
        "abstract tech patterns, particle effects at corners.\n"
        "Pure visual decoration only — absolutely no text anywhere."
    ),

    # ═══════ 第一章分隔页 ═══════
    ("ch1_section", "section", STYLE + STYLE_DARK +
        "Design a chapter section divider slide.\n"
        "Leave top-left corner clean for logo overlay.\n"
        "VERY LARGE number '01' on the left side (oversized, semi-transparent, decorative).\n"
        "Chapter title: '第一章标题' in large white bold text.\n"
        "Subtitle below: '副标题说明' in smaller light text.\n"
        "Right side: decorative geometric shapes (arcs, dots grid)."
    ),

    # ===== 第一章 内容页 =====
    ("ch1_p1", "content", STYLE +
        "Title: '内容页标题'\n"
        "Layout: 左右两栏。\n"
        "左栏：核心要点（3-4 条 bullet points）。\n"
        "右栏：配图区或数据图表区。\n"
        "Bottom right: small page number '03'."
    ),
    ("ch1_p2", "content", STYLE +
        "Title: '内容页标题'\n"
        "Layout: 表格或流程图。\n"
        "4 行 3 列的结构化数据展示。\n"
        "Bottom right: small page number '04'."
    ),

    # ═══════ 第二章分隔页 ═══════
    ("ch2_section", "section", STYLE + STYLE_DARK +
        "Design a chapter section divider slide.\n"
        "Leave top-left corner clean for logo overlay.\n"
        "VERY LARGE number '02' on the left side (oversized, semi-transparent, decorative).\n"
        "Chapter title: '第二章标题' in large white bold text.\n"
        "Subtitle below: '副标题说明' in smaller light text.\n"
        "Right side: decorative geometric shapes."
    ),

    # ===== 第二章 内容页 =====
    ("ch2_p1", "content", STYLE +
        "Title: '内容页标题'\n"
        "Layout: 指标卡片（3-4 个 KPI 数字卡）。\n"
        "每个卡片包含：大号数字 + 指标名称 + 趋势箭头。\n"
        "Bottom right: small page number '06'."
    ),

    # ═══════ 结尾页（纯装饰背景，文字由 python-pptx 叠加） ═══════
    ("closing", "closing", STYLE + STYLE_DARK +
        "Design a PURELY DECORATIVE background for a closing/thank-you slide.\n"
        "STRICTLY NO TEXT, NO CHARACTERS, NO NUMBERS, NO LETTERS, NO WORDS of any language.\n"
        "Leave top-center area clean for logo overlay.\n"
        "Leave the center area clean for text to be overlaid later by code.\n"
        "Decorative: subtle ring or arc pattern, elegant minimal geometric elements, "
        "soft glow effect at center. Pure visual decoration only — absolutely no text anywhere."
    ),
]

# 组装顺序（自动提取）
SLIDE_ORDER = [p[0] for p in POSTERS]

# page_type 查找表
PAGE_TYPES = {p[0]: p[1] for p in POSTERS}

# ╔══════════════════════════════════════════════════════════════╗
# ║  Logo 叠加工具函数                                              ║
# ╚══════════════════════════════════════════════════════════════╝


def _prepare_logo_variants(logo_path):
    """预处理 logo：为深蓝背景生成白色版本。

    原始 logo 是黑色文字 + 绿色云（透明背景），
    深蓝页需要白色文字版本。

    Returns:
        (dark_logo_path, light_logo_path)
        dark  = 黑色文字版（用于浅色背景内容页）
        light = 白色文字版（用于深蓝背景封面/章节/结尾页）
    """
    if not os.path.exists(logo_path):
        print(f"  [WARN] Logo not found: {logo_path}")
        return None, None

    dark_path = logo_path  # 原始就是黑色版

    # 生成白色版：把黑色像素 (#000000) 替换为白色 (#FFFFFF)，绿色保留
    light_dir = os.path.dirname(logo_path)
    light_path = os.path.join(light_dir, "logo-title-white.png")

    if not os.path.exists(light_path):
        img = Image.open(logo_path).convert("RGBA")
        pixels = img.load()
        w, h = img.size
        for y in range(h):
            for x in range(w):
                r, g, b, a = pixels[x, y]
                if a > 10:
                    # 黑色/深灰像素 → 白色（保留绿色）
                    if r < 50 and g < 50 and b < 50:
                        pixels[x, y] = (255, 255, 255, a)
        img.save(light_path)
        print(f"  [OK] Generated white logo: {light_path}")

    return dark_path, light_path


def _add_logo(slide, page_type, dark_logo, light_logo):
    """在 slide 上叠加真实 logo。

    Args:
        slide: pptx Slide 对象
        page_type: "cover" | "section" | "content" | "closing"
        dark_logo: 黑色文字版 logo 路径（浅色背景用）
        light_logo: 白色文字版 logo 路径（深色背景用）
    """
    preset = LOGO_PRESETS.get(page_type, LOGO_PRESETS["content"])

    # 深蓝背景页用白色 logo，浅色背景页用黑色 logo
    if page_type in ("cover", "section", "closing"):
        logo = light_logo
    else:
        logo = dark_logo

    if not logo or not os.path.exists(logo):
        return

    slide.shapes.add_picture(
        logo,
        left=preset["left"],
        top=preset["top"],
        width=preset["width"],
        # height 自动按比例计算
    )


# ╔══════════════════════════════════════════════════════════════╗
# ║  图片生成逻辑                                                  ║
# ╚══════════════════════════════════════════════════════════════╝


def gen_one(key, prompt, size=IMAGE_SIZE, regen=False):
    """生成单张海报图片。"""
    out = os.path.join(ASSETS, f"{key}.png")

    if os.path.exists(out) and not regen:
        print(f"  [SKIP] {key}.png exists")
        return True

    cmd = [
        sys.executable, GEN_SCRIPT,
        "-p", prompt,
        "-s", size,
        "-q", IMAGE_QUALITY,
        "-o", f"{key}.png",
        "-d", ASSETS,
    ]

    env = os.environ.copy()
    env["PYTHONIOENCODING"] = "utf-8"

    try:
        result = subprocess.run(
            cmd,
            stdout=subprocess.PIPE,
            stderr=subprocess.PIPE,
            timeout=TIMEOUT_SEC,
            env=env,
        )
        stdout = result.stdout.decode("utf-8", errors="replace") if result.stdout else ""
        stderr = result.stderr.decode("utf-8", errors="replace") if result.stderr else ""

        if result.returncode == 0 and os.path.exists(out):
            sz = os.path.getsize(out)
            print(f"  [OK] {key}.png ({sz // 1024} KB)")
            return True
        else:
            print(f"  [FAIL] {key} rc={result.returncode}")
            if stderr:
                print(f"    stderr: {stderr[:300]}")
            return False

    except subprocess.TimeoutExpired:
        print(f"  [TIMEOUT] {key} (>{TIMEOUT_SEC}s)")
        return False
    except Exception as e:
        print(f"  [ERROR] {key}: {e}")
        return False


def gen_posters(only=None, regen=False):
    """批量生成海报图片。"""
    targets = POSTERS
    if only:
        targets = [p for p in POSTERS if p[0] == only]
        if not targets:
            print(f"[ERROR] Unknown key: {only}")
            print(f"  Available: {[p[0] for p in POSTERS]}")
            return

    total = len(targets)
    print(f"Generating {total} poster image(s)...")
    print(f"  Output: {ASSETS}")
    print(f"  Size: {IMAGE_SIZE} | Quality: {IMAGE_QUALITY}")
    print()

    ok, fail = 0, 0
    for i, (key, page_type, prompt) in enumerate(targets, 1):
        label = {"cover": "封面", "section": "章节分隔",
                 "content": "内容", "closing": "结尾"}.get(page_type, page_type)
        print(f"[{i}/{total}] {key} ({label})")
        if gen_one(key, prompt, regen=regen):
            ok += 1
        else:
            fail += 1
        if i < total:
            time.sleep(SLEEP_SEC)

    print(f"\n{'=' * 50}")
    print(f"Done: {ok} OK, {fail} FAIL")
    if fail > 0:
        print("Re-run with --regen to retry failed items.")


# ╔══════════════════════════════════════════════════════════════╗
# ║  PPT 组装逻辑（含 logo 叠加）                                   ║
# ╚══════════════════════════════════════════════════════════════╝


def assemble_ppt(output_name="output_poster.pptx"):
    """将海报图片组装为 PPT，并叠加真实 logo。"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank_layout = prs.slide_layouts[6]

    # 预处理 logo（生成白色版本）
    dark_logo, light_logo = _prepare_logo_variants(LOGO_PATH)
    if dark_logo:
        print(f"  Logo (dark bg): {light_logo}")
        print(f"  Logo (light bg): {dark_logo}")
    else:
        print("  [WARN] No logo found, slides will have no logo overlay.")

    added, skipped = 0, 0
    print(f"Assembling PPT from {len(SLIDE_ORDER)} slides...")

    for key in SLIDE_ORDER:
        img_path = os.path.join(ASSETS, f"{key}.png")
        if not os.path.exists(img_path):
            print(f"  [WARN] Missing: {key}.png - skipped")
            skipped += 1
            continue

        slide = prs.slides.add_slide(blank_layout)

        # 1. 先插入全屏海报底图
        slide.shapes.add_picture(
            img_path,
            left=Emu(0),
            top=Emu(0),
            width=SLIDE_WIDTH,
            height=SLIDE_HEIGHT,
        )

        # 2. 再叠加真实 logo（在底图上方）
        page_type = PAGE_TYPES.get(key, "content")
        _add_logo(slide, page_type, dark_logo, light_logo)

        added += 1

    out_path = os.path.join(BASE_DIR, output_name)
    prs.save(out_path)
    print(f"\nSaved: {out_path}")
    print(f"  Slides: {added} added, {skipped} skipped")
    return out_path


# ╔══════════════════════════════════════════════════════════════╗
# ║  命令行入口                                                    ║
# ╚══════════════════════════════════════════════════════════════╝


def main():
    parser = argparse.ArgumentParser(description="360 PPT Poster Generator (v3)")
    parser.add_argument("action", choices=["gen", "build", "all"],
                        help="gen=生成图片, build=组装PPT, all=全流程")
    parser.add_argument("--only", type=str, default=None,
                        help="仅生成指定 key 的图片")
    parser.add_argument("--regen", action="store_true",
                        help="强制重生成（覆盖已有图片）")
    parser.add_argument("--output", type=str, default="output_poster.pptx",
                        help="PPT 输出文件名")

    args = parser.parse_args()

    if args.action in ("gen", "all"):
        gen_posters(only=args.only, regen=args.regen)

    if args.action in ("build", "all"):
        assemble_ppt(output_name=args.output)


if __name__ == "__main__":
    main()
