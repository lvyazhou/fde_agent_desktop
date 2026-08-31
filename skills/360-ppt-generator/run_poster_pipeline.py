# -*- coding: utf-8 -*-
"""
360 PPT 海报模式 — 平台在线生成入口脚本
======================================

子 agent 只需要：
1. 在 shared/{topic}/posters.json 写入页面列表
2. 执行本脚本：python run_poster_pipeline.py --workspace WORKSPACE_ROOT --topic TOPIC

本脚本自动完成：
- 读取 posters.json
- 用 ThreadPoolExecutor(max_workers=6) 并行生成所有海报图片
- 验证所有图片存在且 >10KB
- 组装为 PPT（含 Logo 叠加 + 封面/结尾可编辑文字）

posters.json 格式：
[
  {"key": "cover", "page_type": "cover", "prompt": "..."},
  {"key": "toc", "page_type": "content", "prompt": "..."},
  {"key": "section_1", "page_type": "section", "prompt": "..."},
  ...
  {"key": "closing", "page_type": "closing", "prompt": "..."}
]

page_type 取值: "cover" | "section" | "content" | "closing"
"""

import os
import sys
import json
import time
import subprocess
import argparse
from concurrent.futures import ThreadPoolExecutor, as_completed

# 强制 stdout 行缓冲——确保 pipe 模式下 print() 逐行刷出（供 ExecTool 实时读取）
sys.stdout.reconfigure(line_buffering=True)

# ============ 路径（由命令行参数传入，不自行计算） ============

def parse_args():
    parser = argparse.ArgumentParser(description="360 PPT 海报批量生成+组装")
    parser.add_argument("--workspace", required=True, help="workspace 根目录绝对路径")
    parser.add_argument("--topic", required=True, help="主题名称（英文/拼音）")
    parser.add_argument("--action", default="all", choices=["gen", "build", "all"],
                        help="gen=生成图片, build=组装PPT, all=全流程")
    parser.add_argument("--workers", type=int, default=6, help="并行线程数")
    parser.add_argument("--theme", default="blue", choices=["blue", "green"],
                        help="配色主题：blue=360深蓝(默认), green=360绿色系(教育/培训版)")
    return parser.parse_args()


def main():
    args = parse_args()

    WORKSPACE_ROOT = args.workspace
    TOPIC = args.topic
    PROJECT_DIR = os.path.join(WORKSPACE_ROOT, "shared", TOPIC)
    ASSETS_DIR = os.path.join(PROJECT_DIR, "assets")
    # 优先查找 .claude/skills，其次 skills/
    SKILL_DIR = os.path.join(WORKSPACE_ROOT, ".claude", "skills", "360-ppt-generator")
    if not os.path.isdir(SKILL_DIR):
        SKILL_DIR = os.path.join(WORKSPACE_ROOT, "skills", "360-ppt-generator")
    GEN_SCRIPT = os.path.join(SKILL_DIR, "generate_images.py")
    BRAND_DIR = os.path.join(SKILL_DIR, "brand_assets")
    OUTPUT_PPTX = os.path.join(PROJECT_DIR, f"{TOPIC}.pptx")
    POSTERS_JSON = os.path.join(PROJECT_DIR, "posters.json")

    os.makedirs(ASSETS_DIR, exist_ok=True)

    # 验证必要文件
    if not os.path.exists(GEN_SCRIPT):
        print(f"[FAIL] generate_images.py not found: {GEN_SCRIPT}")
        sys.exit(1)
    if not os.path.exists(POSTERS_JSON):
        print(f"[FAIL] posters.json not found: {POSTERS_JSON}")
        print("子 agent 需要先创建 posters.json 文件，格式见脚本头部注释。")
        sys.exit(1)

    # 读取海报列表
    with open(POSTERS_JSON, "r", encoding="utf-8") as f:
        posters = json.load(f)
    print(f"Loaded {len(posters)} posters from {POSTERS_JSON}")

    # ============ STYLE 定义 ============
    STYLE_LIGHT = (
        "Professional Chinese business presentation slide poster, "
        "CLEAN LIGHT BACKGROUND (very light blue-gray #F8F9FD) "
        "with subtle blue geometric corner decorations "
        "(top-right and bottom-left corners), "
        "DO NOT draw any logo or brand mark anywhere on the image. "
        "Leave the top-left corner area (top 0.8 inches, left 2 inches) completely clean "
        "with no text, icons, or decorations — a real logo will be overlaid later. "
        "Modern flat design, corporate cybersecurity tech style, "
        "title text in deep blue (#0823AB) bold font, "
        "body text in dark gray (#333333) with light blue-white (#E5E9FA) info-card backgrounds, "
        "accent highlights in green (#00C853) and orange (#FF6D00) for warnings, "
        "clear Chinese text typography with bold titles and organized content, "
        "16:9 widescreen aspect ratio. "
        "The text must be rendered clearly and accurately in Simplified Chinese. "
    )

    STYLE_DARK = (
        "OVERRIDE STYLE: Use a DEEP BLUE gradient background (#0B1D4A to #0823AB) "
        "instead of light background. "
        "DO NOT draw any logo or brand mark — a real logo will be overlaid later. "
        "All text in WHITE or light blue (#E5E9FA). "
        "Decorative: subtle tech grid pattern, circuit board traces, geometric accents, "
        "glowing particle effects. "
    )

    # ============ 绿色系 STYLE（教育/培训版，白底+绿色装饰） ============
    STYLE_GREEN_LIGHT = (
        "Professional Chinese business training-slide poster, clean modern corporate style. "
        "CLEAN WHITE BACKGROUND with a subtle light-green tint, airy with generous whitespace. "
        "PRIMARY BRAND COLOR is GREEN (#00A870 / #16B37F emerald green) used for the title bar accent, "
        "icons, numbered badges, progress arrows and highlights; use dark near-black text for body. "
        "A thin decorative tech-line pattern faint in the top-right corner only. "
        "CRITICAL: leave the TOP-LEFT corner (top 1.1 inch, left 3 inch) COMPLETELY EMPTY for a logo overlay - "
        "no text, no graphics there. "
        "STRICTLY NO LOGO, NO BRAND MARK, NO '360' TEXT anywhere - a real logo PNG will be overlaid later. "
        "Put the page title at the TOP-CENTER in bold, clear of the top-left logo zone. "
        "Layout style: equal columns/modules laid out horizontally across the middle, each module "
        "has a small green numbered badge, a small green line-icon, a bold heading, and a few short bullet lines. "
        "A horizontal SUMMARY BAR at the bottom with several equal cells (green background or green icons). "
        "Modern flat design, rounded cards with very soft shadows, thin green dividers between columns. "
        "FONT: render ALL Chinese text in a clean modern sans-serif GOTHIC typeface like Microsoft YaHei / Source Han Sans - "
        "no serif, no handwriting fonts, consistent across the whole slide. "
        "FONT-SIZE HIERARCHY: page title largest and bold (about 40pt), module headings medium bold (about 22-24pt), "
        "body bullets regular (about 16-18pt), captions smaller (about 13-14pt). "
        "16:9 widescreen. The text must be rendered clearly and accurately in Simplified Chinese. "
        "Match the SAME visual style across all slides in this set: same green, same title position, same column layout, same bottom bar. "
    )

    STYLE_GREEN_DARK = (
        "OVERRIDE STYLE: deep emerald-green gradient background (#0A5F44 to #00A870) "
        "instead of white background. "
        "STRICTLY NO LOGO, NO BRAND MARK, NO '360' TEXT anywhere - a real logo PNG will be overlaid later. "
        "Leave the top-center/top-left area clean for logo overlay. "
        "All text in WHITE or light green (#D4F5E9). "
        "Decorative: subtle tech-line pattern, geometric accents. "
        "This is a COVER/SECTION DIVIDER/CLOSING slide - use deep green gradient NOT white background. "
    )

    def get_style(page_type):
        if args.theme == "green":
            if page_type in ("cover", "section", "closing"):
                return STYLE_GREEN_LIGHT + STYLE_GREEN_DARK
            return STYLE_GREEN_LIGHT
        # 默认蓝色系
        if page_type in ("cover", "section", "closing"):
            return STYLE_LIGHT + STYLE_DARK
        return STYLE_LIGHT

    # ============ 绿色海报版页型模板 ============
    # 这些是给 posters.json 写 prompt 时的默认口径。
    #
    # 普通图文页：
    # - 标题 + 主视觉 + 3~4 条短说明 + 底部结论条
    # - 适合讲为什么、场景感、痛点、扶持说明、客户故事
    #
    # 结构化图文页：
    # - 标题 + 流程 / 对比 / 矩阵 / 关系图 + 底部结论条
    # - 适合讲怎么做、关系怎么连、流程怎么走
    #
    # 默认混排顺序：封面 → 目录 → 章节页 → 普通图文页 → 结构化图文页 → 普通图文页 → 结构化图文页 → 结尾页

    # ============ Step 1: 生成图片 ============
    if args.action in ("gen", "all"):
        print(f"\n{'='*60}")
        print(f"Step 1: 生成 {len(posters)} 张海报图片 (max_workers={args.workers})")
        print(f"{'='*60}")

        def gen_one(poster):
            key = poster["key"]
            page_type = poster.get("page_type", "content")
            prompt = poster["prompt"]
            out = os.path.join(ASSETS_DIR, f"{key}.png")

            if os.path.exists(out) and os.path.getsize(out) > 10240:
                return (key, True, f"[SKIP] {key}.png exists ({os.path.getsize(out)//1024}KB)")

            full_prompt = get_style(page_type) + "\n\n" + prompt
            t0 = time.time()
            env = {**os.environ, "PYTHONIOENCODING": "utf-8"}
            try:
                result = subprocess.run(
                    [sys.executable, GEN_SCRIPT,
                     "-p", full_prompt, "-s", "1792x1024", "-q", "high",
                     "-o", f"{key}.png", "-d", ASSETS_DIR],
                    stdout=subprocess.PIPE, stderr=subprocess.PIPE,
                    timeout=300, env=env
                )
                elapsed = time.time() - t0
                if result.returncode == 0 and os.path.exists(out) and os.path.getsize(out) > 10240:
                    return (key, True, f"[OK] {key}.png ({os.path.getsize(out)//1024}KB, {elapsed:.1f}s)")
                stderr = result.stderr.decode("utf-8", errors="replace")[:200] if result.stderr else ""
                return (key, False, f"[FAIL] {key} rc={result.returncode} {stderr}")
            except subprocess.TimeoutExpired:
                return (key, False, f"[TIMEOUT] {key} (>300s)")
            except Exception as e:
                return (key, False, f"[ERROR] {key}: {e}")

        t_start = time.time()
        results = {}
        total = len(posters)
        done_count = 0
        with ThreadPoolExecutor(max_workers=args.workers) as pool:
            futures = {pool.submit(gen_one, p): p["key"] for p in posters}
            for f in as_completed(futures):
                k = futures[f]
                results[k] = f.result()
                done_count += 1
                # 标准进度格式: [PROGRESS] current/total | detail
                # 前端 ChatRightPanel 解析此格式展示子进度条
                print(f"[PROGRESS] {done_count}/{total} | {results[k][2]}")

        ok = sum(1 for v in results.values() if v[1])
        print(f"\n[PROGRESS] {total}/{total} | Done: {ok}/{total} OK ({time.time()-t_start:.1f}s)")

        if ok < total:
            failed = [k for k, v in results.items() if not v[1]]
            print(f"[WARN] Failed: {failed}")

    # ============ Step 2: 验证 ============
    print(f"\n{'='*60}")
    print("Step 2: 验证图片完整性")
    print(f"{'='*60}")

    missing = []
    for p in posters:
        path = os.path.join(ASSETS_DIR, f"{p['key']}.png")
        if not os.path.exists(path):
            missing.append(p["key"])
        elif os.path.getsize(path) < 10240:
            missing.append(f"{p['key']} (too small)")

    if missing:
        print(f"[PROGRESS] {len(posters)-len(missing)}/{len(posters)} | 验证失败: 缺少 {len(missing)} 张图片")
        if args.action == "all":
            print("Skipping PPT assembly due to missing images.")
            sys.exit(1)
    else:
        print(f"[PROGRESS] {len(posters)}/{len(posters)} | 验证通过: 全部 {len(posters)} 张图片完整")

    # ============ Step 3: 组装 PPT ============
    if args.action in ("build", "all"):
        print(f"\n{'='*60}")
        print("Step 3: 组装 PPT")
        print(f"{'='*60}")

        from pptx import Presentation
        from pptx.util import Inches, Emu, Pt
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN

        SLIDE_W = Inches(13.333)
        SLIDE_H = Inches(7.5)

        # Logo 配置
        LOGO_ORIG = os.path.join(BRAND_DIR, "logo-title.png")
        LOGO_WHITE = os.path.join(BRAND_DIR, "logo-title-white.png")

        LOGO_PRESETS = {
            "cover":   {"logo": LOGO_WHITE, "left": 5.2, "top": 0.5, "width": 3.0},
            "section": {"logo": LOGO_WHITE, "left": 0.5, "top": 0.3, "width": 2.2},
            "content": {"logo": LOGO_ORIG,  "left": 0.2, "top": 0.1, "width": 2.2},
            "closing": {"logo": LOGO_WHITE, "left": 5.2, "top": 0.8, "width": 3.0},
        }

        prs = Presentation()
        prs.slide_width = SLIDE_W
        prs.slide_height = SLIDE_H
        blank_layout = prs.slide_layouts[6]

        for p in posters:
            key = p["key"]
            page_type = p.get("page_type", "content")
            title = p.get("title", "")
            subtitle = p.get("subtitle", "")
            img_path = os.path.join(ASSETS_DIR, f"{key}.png")

            if not os.path.exists(img_path):
                print(f"  [SKIP] {key}.png not found")
                continue

            slide = prs.slides.add_slide(blank_layout)
            # 全屏背景图
            slide.shapes.add_picture(img_path, Emu(0), Emu(0), SLIDE_W, SLIDE_H)

            # Logo 叠加
            preset = LOGO_PRESETS.get(page_type, LOGO_PRESETS["content"])
            logo_file = preset["logo"]
            if os.path.exists(logo_file):
                slide.shapes.add_picture(
                    logo_file,
                    Inches(preset["left"]), Inches(preset["top"]),
                    Inches(preset["width"])
                )

            # 封面页/结尾页：叠加可编辑文字（双击可改）
            if page_type == "cover" and title:
                # 标题：白色大字居中
                txBox = slide.shapes.add_textbox(
                    Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                para.text = title
                para.font.size = Pt(44)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para.alignment = PP_ALIGN.CENTER
                # 副标题
                if subtitle:
                    sub_box = slide.shapes.add_textbox(
                        Inches(2.0), Inches(4.2), Inches(9.3), Inches(1.0)
                    )
                    stf = sub_box.text_frame
                    stf.word_wrap = True
                    sp = stf.paragraphs[0]
                    sp.text = subtitle
                    sp.font.size = Pt(22)
                    sp.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
                    sp.alignment = PP_ALIGN.CENTER

            elif page_type == "closing" and title:
                # 结尾页：感谢文字居中
                txBox = slide.shapes.add_textbox(
                    Inches(2.0), Inches(2.8), Inches(9.3), Inches(1.5)
                )
                tf = txBox.text_frame
                tf.word_wrap = True
                para = tf.paragraphs[0]
                para.text = title
                para.font.size = Pt(40)
                para.font.bold = True
                para.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                para.alignment = PP_ALIGN.CENTER
                if subtitle:
                    sub_box = slide.shapes.add_textbox(
                        Inches(3.0), Inches(4.5), Inches(7.3), Inches(0.8)
                    )
                    stf = sub_box.text_frame
                    sp = stf.paragraphs[0]
                    sp.text = subtitle
                    sp.font.size = Pt(18)
                    sp.font.color.rgb = RGBColor(0x90, 0xCA, 0xF9)
                    sp.alignment = PP_ALIGN.CENTER

            print(f"  [+] {key} ({page_type})")
            # 输出组装进度
            slide_idx = posters.index(p) + 1
            print(f"[PROGRESS] {slide_idx}/{len(posters)} | 组装幻灯片: {key} ({page_type})")

        prs.save(OUTPUT_PPTX)
        pptx_size_kb = os.path.getsize(OUTPUT_PPTX) // 1024
        print(f"\n[PROGRESS] {len(prs.slides)}/{len(prs.slides)} | PPT 生成完成: {len(prs.slides)} 页, {pptx_size_kb}KB")
        print(f"[OK] PPT saved: {OUTPUT_PPTX}")


if __name__ == "__main__":
    main()
